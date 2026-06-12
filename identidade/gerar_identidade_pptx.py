from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
from lxml import etree
import copy, os

ID_DIR = "/Users/jessicagregorio/Oca Propaganda/identidade"
LOGO_PRINCIPAL = os.path.join(ID_DIR, "logo_pptx_principal.png")
LOGO_AMARELO   = os.path.join(ID_DIR, "logo_pptx_amarelo.png")
LOGO_CLARO     = os.path.join(ID_DIR, "logo_pptx_claro.png")

def add_logo(slide, path, x, y, w, h):
    slide.shapes.add_picture(path, x, y, w, h)

prs = Presentation()
prs.slide_width  = Inches(16)
prs.slide_height = Inches(9)
W = prs.slide_width
H = prs.slide_height
BLANK = prs.slide_layouts[6]

# Colors
Y   = RGBColor(0xFF, 0xDE, 0x21)
BK  = RGBColor(0x0A, 0x0A, 0x0A)
OW  = RGBColor(0xF5, 0xF5, 0xF0)
MG  = RGBColor(0x80, 0x80, 0x80)
LG  = RGBColor(0x9B, 0x9B, 0x9B)
WH  = RGBColor(0xFF, 0xFF, 0xFF)
D16 = RGBColor(0x16, 0x16, 0x16)
D20 = RGBColor(0x20, 0x20, 0x20)
GBG = RGBColor(0xEA, 0xF3, 0xDE)
GTX = RGBColor(0x27, 0x50, 0x0A)
RBG = RGBColor(0xFC, 0xEB, 0xEB)
RTX = RGBColor(0xA3, 0x2D, 0x2D)

def i(v): return Inches(v)

def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, fill_c, line_c=None, line_w=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_c
    if line_c:
        s.line.color.rgb = line_c
        s.line.width = Pt(line_w or 0.75)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, x, y, w, h, size, bold=False, color=BK,
        align=PP_ALIGN.LEFT, italic=False, font_name="Arial"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    # remove internal padding
    tb.left = x; tb.top = y; tb.width = w; tb.height = h
    return tb

def txt2(slide, lines, x, y, w, h, size, bold=False, color=BK,
         align=PP_ALIGN.LEFT, line_spacing_pt=None):
    """Multi-line text with consistent formatting."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = "Arial"
    return tb

# ═══════════════════════════════════════════════════════
# SLIDE 1 — LOGOTIPO PRINCIPAL
# ═══════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, BK)

# barra amarela esquerda
rect(s1, 0, 0, i(0.28), H, Y)

# Logo oficial: PROPAGANDA (Cactus) + OCA (Oferta do Dia)
add_logo(s1, LOGO_PRINCIPAL, i(0.5), i(1.1), i(10), i(4.0))

# label técnico
txt(s1, "// IDENTIDADE VISUAL  ·  2025", i(0.5), i(5.25),
    i(8), i(0.35), size=10, color=MG)

# tagline
txt(s1, "Marketing que funciona de verdade.", i(0.5), i(5.65),
    i(8), i(0.5), size=16, color=LG)

# círculo geométrico — canto superior direito
circ = s1.shapes.add_shape(9, i(13.5), i(0.3), i(2.2), i(2.2))
circ.fill.background()
circ.line.color.rgb = Y
circ.line.width = Pt(1.5)

circ2 = s1.shapes.add_shape(9, i(13.7), i(0.5), i(1.8), i(1.8))
circ2.fill.background()
circ2.line.color.rgb = RGBColor(0x35, 0x35, 0x35)
circ2.line.width = Pt(0.75)

# ponto central
dot = s1.shapes.add_shape(9, i(14.52), i(1.32), i(0.06), i(0.06))
dot.fill.solid(); dot.fill.fore_color.rgb = Y
dot.line.fill.background()

# rodapé amarelo
rect(s1, 0, i(8.35), W, i(0.65), Y)
foot_items = ["propagandaoca.com.br", "oi@propagandaoca.com.br",
              "(11) 9 5912-5935", "Bom Jesus dos Perdoes  .  SP"]
slot = int(i(16) / len(foot_items))
for idx, item in enumerate(foot_items):
    txt(s1, item, idx * slot, i(8.52), slot, i(0.35),
        size=9, color=BK, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════
# SLIDE 2 — PALETA DE CORES
# ═══════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, OW)

# header
rect(s2, 0, 0, W, i(0.9), BK)
txt(s2, "IDENTIDADE VISUAL", i(0.5), i(0.1), i(9), i(0.7),
    size=30, bold=True, color=Y)
txt(s2, "02 / CORES", i(13.5), i(0.28), i(2.3), i(0.4),
    size=9, color=MG)

cores = [
    (Y,   "#FFDE21", "255 . 222 . 33",  "0 . 13 . 87 . 0",  "AMARELO",
     "Cor de identidade.\nDestaques e CTAs.", BK),
    (BK,  "#0A0A0A", "10 . 10 . 10",    "0 . 0 . 0 . 96",   "PRETO",
     "Fundos principais.\nTexto sobre fundo claro.", Y),
    (OW,  "#F5F5F0", "245 . 245 . 240",  "0 . 0 . 2 . 4",    "OFF-WHITE",
     "Fundos secundarios.\nPapelaria e apresentacoes.", BK),
    (RGBColor(0x1A,0x1A,0x1A), "#1A1A1A", "26 . 26 . 26", "0 . 0 . 0 . 90",
     "CINZA ESCURO", "Texto e elementos\nde suporte.", LG),
]
sw = i(3.6); sh = i(4.5); pad = i(0.4); gap = i(0.27)
sy = i(1.1)
for idx, (col, hex_, rgb_, cmyk_, name, uso, tc) in enumerate(cores):
    sx = pad + idx*(sw+gap)
    r = rect(s2, sx, sy, sw, sh, col)
    if col == OW:
        r.line.color.rgb = RGBColor(0xC8,0xC8,0xC8); r.line.width = Pt(0.5)
    txt(s2, name, sx+i(0.15), sy+i(0.18), sw-i(0.3), i(0.55),
        size=20, bold=True, color=tc)
    txt(s2, hex_, sx+i(0.15), sy+i(0.72), sw-i(0.3), i(0.4),
        size=13, bold=True, color=tc if col != OW else MG)
    # dados abaixo
    dy = sy+sh+i(0.18)
    txt(s2, "RGB",  sx, dy,         sw, i(0.25), size=9, color=LG)
    txt(s2, rgb_,   sx, dy+i(0.22), sw, i(0.28), size=10, color=BK)
    txt(s2, "CMYK", sx, dy+i(0.55), sw, i(0.25), size=9, color=LG)
    txt(s2, cmyk_,  sx, dy+i(0.77), sw, i(0.28), size=10, color=BK)
    txt(s2, uso,    sx, dy+i(1.15), sw, i(0.55), size=10, color=MG)

# barra inferior nos swatches
for idx, (col, *_) in enumerate(cores):
    sx = pad + idx*(sw+gap)
    bar_c = Y if col != Y else BK
    rect(s2, sx, sy+sh-i(0.06), sw, i(0.06), bar_c)

# nota rodapé
txt(s2, "PROPAGANDA OCA  —  BRANDBOOK 2025", i(0.5), i(8.78),
    i(8), i(0.2), size=8, color=LG)

# ═══════════════════════════════════════════════════════
# SLIDE 3 — TIPOGRAFIA
# ═══════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, BK)

rect(s3, 0, 0, W, i(0.9), Y)
txt(s3, "TIPOGRAFIA", i(0.5), i(0.1), i(9), i(0.7),
    size=30, bold=True, color=BK)
txt(s3, "03 / TYPE", i(13.8), i(0.28), i(2), i(0.4),
    size=9, color=BK)

fontes = [
    ("LOGOTIPO / PROPAGANDA", "Cactus Classical", 36, Y,  "Serifada brasileira. Para o\nnome PROPAGANDA no logo."),
    ("LOGOTIPO / OCA",        "Oferta do Dia",    52, OW, "Display bold nacional.\nPara OCA no logo da marca."),
    ("CORPO E DADOS",         "DM  Mono",         22, LG, "Precisao, tecnica, referencia.\nPara numeros e labels."),
]
bw = i(4.9); bh = i(3.1); by = i(1.05); bpad = i(0.35); gap3 = i(0.15)
for idx, (role, sample, ssize, tc, desc) in enumerate(fontes):
    bx = bpad + idx*(bw+gap3)
    rect(s3, bx, by, bw, bh, D20)
    txt(s3, role,   bx+i(0.18), by+i(0.18), bw-i(0.3), i(0.3),
        size=8, color=Y)
    txt(s3, sample, bx+i(0.18), by+i(0.52), bw-i(0.3), i(1.2),
        size=ssize, bold=True, color=tc)
    for li, line in enumerate(desc.split("\n")):
        txt(s3, line, bx+i(0.18), by+bh-i(0.78)+li*i(0.38),
            bw-i(0.3), i(0.35), size=10, color=MG)

# faixa alfabeto
rect(s3, 0, i(4.25), W, i(1.05), Y)
txt(s3, "ABCDEFGHIJKLMNOPQRSTUVWXYZ",
    i(0.18), i(4.28), i(15.6), i(0.95),
    size=54, bold=True, color=BK)

# faixa numérica
rect(s3, 0, i(5.38), W, i(1.0), D16)
txt(s3, "0123456789   . , % ; x",
    i(0.22), i(5.42), i(14), i(0.85),
    size=46, color=Y)

# pairing guide
rect(s3, 0, i(6.46), W, i(2.54), RGBColor(0x0E,0x0E,0x0E))
txt(s3, "PAIRING GUIDE", i(0.4), i(6.62), i(4), i(0.3),
    size=9, color=LG)
pairs = [
    ("Cactus Classical Serif", "PROPAGANDA no logo. Serifada com identidade brasileira."),
    ("Oferta do Dia", "OCA no logo. Display bold nacional de alto impacto."),
    ("DM Mono Regular", "Numeros, metricas, labels tecnicos, URLs, datas"),
]
for ji, (pname, pdesc) in enumerate(pairs):
    py3 = i(7.0) + ji*i(0.54)
    txt(s3, pname, i(0.4), py3, i(4.5), i(0.42),
        size=12, bold=True, color=WH)
    rect(s3, i(5.1), py3+i(0.15), i(0.22), i(0.08), Y)
    txt(s3, pdesc, i(5.45), py3, i(9), i(0.42),
        size=12, color=MG)

# ═══════════════════════════════════════════════════════
# SLIDE 4 — MANIFESTO VISUAL
# ═══════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, BK)

# grid de linhas finas decorativas
for xi in range(0, 17):
    lx = i(xi)
    r2 = rect(s4, lx, 0, i(0.005), H, RGBColor(0x16,0x16,0x16))

# bloco amarelo — quadrante superior direito
rect(s4, i(8.3), 0, i(7.7), i(4.55), Y)

# círculo geométrico
circ3 = s4.shapes.add_shape(9, i(9.8), i(0.4), i(3.1), i(3.1))
circ3.fill.solid(); circ3.fill.fore_color.rgb = BK
circ3.line.fill.background()
circ4 = s4.shapes.add_shape(9, i(10.05), i(0.65), i(2.6), i(2.6))
circ4.fill.background()
circ4.line.color.rgb = Y; circ4.line.width = Pt(1.5)
dot2 = s4.shapes.add_shape(9, i(11.3), i(1.9), i(0.07), i(0.07))
dot2.fill.solid(); dot2.fill.fore_color.rgb = Y
dot2.line.fill.background()

# "PROPAGANDA" vertical — letras empilhadas
chars_p = list("PROPAGANDA")
char_h = i(0.82)
start_v = int((H - len(chars_p)*char_h) / 2)
for ci, ch in enumerate(chars_p):
    txt(s4, ch, i(0.05), start_v + ci*char_h, i(0.7), char_h,
        size=44, bold=True, color=Y)

# "OCA" enorme
txt(s4, "OCA", i(0.88), i(4.4), i(7.2), i(3.2),
    size=160, bold=True, color=WH)

# linha acima de OCA
rect(s4, i(0.88), i(4.35), i(7.3), i(0.03), Y)

# manifesto — metade inferior direita
txt(s4, "Marketing que funciona de verdade.",
    i(8.4), i(4.7), i(7.2), i(0.6),
    size=18, bold=True, color=WH)
txt(s4, "Feito no interior. Para o interior.",
    i(8.4), i(5.35), i(7.2), i(0.5),
    size=16, color=LG)
txt(s4, "Sem jargao. Sem template. Sem enrolacao.",
    i(8.4), i(5.85), i(7.2), i(0.5),
    size=16, color=LG)

# bloco métricas
rect(s4, 0, i(7.2), i(8.5), H-i(7.2), D16)
metrics = [("+120","EMPRESAS"), ("3x","MAIS LEADS"), ("87%","RETENCAO"), ("5","ANOS")]
mslot = int(i(8.5) / len(metrics))
for mi, (num, lbl) in enumerate(metrics):
    mx4 = i(0.55) + mi*mslot
    txt(s4, num, mx4, i(7.25), mslot-i(0.2), i(1.3),
        size=58, bold=True, color=Y)
    txt(s4, lbl, mx4, i(8.55), mslot-i(0.2), i(0.35),
        size=10, color=LG)

# linha inferior amarela
rect(s4, 0, i(8.97), W, i(0.03), Y)
txt(s4, "PROPAGANDA OCA  —  2025",
    i(12.5), i(8.7), i(3.3), i(0.25), size=8, color=MG)

# ═══════════════════════════════════════════════════════
# SLIDE 5 — SISTEMA DE APLICACAO
# ═══════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5, OW)

rect(s5, 0, 0, W, i(0.9), BK)
txt(s5, "SISTEMA DE APLICACAO", i(0.5), i(0.1), i(9), i(0.7),
    size=30, bold=True, color=Y)
txt(s5, "05 / GRID", i(13.8), i(0.28), i(2), i(0.4),
    size=9, color=MG)

# Cartão de visita — FRENTE
cw5, ch5 = i(5.6), i(3.2)
cx5, cy5 = i(0.4), i(1.05)
rect(s5, cx5, cy5, cw5, ch5, BK)
rect(s5, cx5, cy5, i(0.2), ch5, Y)                     # barra lateral
add_logo(s5, LOGO_PRINCIPAL, cx5+i(0.25), cy5+i(0.2), cw5-i(0.45), i(1.3))
rect(s5, cx5+i(0.3), cy5+ch5-i(1.12), cw5-i(0.6), i(0.015),
     RGBColor(0x35,0x35,0x35))
txt(s5, "Jessica Gregorio", cx5+i(0.3), cy5+ch5-i(1.1),
    cw5-i(0.4), i(0.38), size=13, bold=True, color=WH)
txt(s5, "Diretora  .  Propaganda Oca", cx5+i(0.3), cy5+ch5-i(0.75),
    cw5-i(0.4), i(0.3), size=10, color=LG)
txt(s5, "(11) 9 5912-5935", cx5+i(0.3), cy5+ch5-i(0.5),
    cw5-i(0.4), i(0.26), size=10, color=LG)
txt(s5, "oi@propagandaoca.com.br", cx5+i(0.3), cy5+ch5-i(0.26),
    cw5-i(0.4), i(0.24), size=10, color=LG)
txt(s5, "CARTAO DE VISITA — FRENTE", cx5, cy5+ch5+i(0.08),
    cw5, i(0.2), size=8, color=MG)

# Cartão de visita — VERSO
vx5 = cx5+cw5+i(0.25)
rect(s5, vx5, cy5, cw5, ch5, Y)
add_logo(s5, LOGO_AMARELO, vx5+i(0.3), cy5+i(0.7), cw5-i(0.5), i(1.3))
txt(s5, "CARTAO DE VISITA — VERSO", vx5, cy5+ch5+i(0.08),
    cw5, i(0.2), size=8, color=MG)

# POST INSTAGRAM 1:1
ig_x5, ig_y5 = i(0.4), i(4.7)
ig_sz = i(3.45)
rect(s5, ig_x5, ig_y5, ig_sz, ig_sz, BK)
txt(s5, "3x",    ig_x5+i(0.2), ig_y5+i(0.25), ig_sz-i(0.3), i(1.1),
    size=66, bold=True, color=Y)
txt(s5, "MAIS",  ig_x5+i(0.2), ig_y5+i(1.25), ig_sz-i(0.3), i(0.9),
    size=60, bold=True, color=WH)
txt(s5, "LEADS", ig_x5+i(0.2), ig_y5+i(2.1),  ig_sz-i(0.3), i(0.9),
    size=60, bold=True, color=WH)
rect(s5, ig_x5, ig_y5+ig_sz-i(0.45), ig_sz, i(0.45), Y)
txt(s5, "@propaganda.oca", ig_x5+i(0.14), ig_y5+ig_sz-i(0.36),
    ig_sz-i(0.2), i(0.3), size=10, color=BK)
txt(s5, "POST INSTAGRAM 1:1", ig_x5, ig_y5+ig_sz+i(0.08),
    ig_sz, i(0.2), size=8, color=MG)

# PAPEL TIMBRADO A4
pt_x5 = ig_x5+ig_sz+i(0.3)
pt_w5, pt_h5 = i(2.4), i(3.4)
pt_y5 = ig_y5
rect(s5, pt_x5, pt_y5, pt_w5, pt_h5, WH,
     line_c=RGBColor(0xC8,0xC8,0xC8), line_w=0.5)
rect(s5, pt_x5, pt_y5, pt_w5, i(0.42), BK)
rect(s5, pt_x5, pt_y5, i(0.09), i(0.42), Y)
add_logo(s5, LOGO_PRINCIPAL, pt_x5+i(0.14), pt_y5+i(0.02), pt_w5-i(0.25), i(0.38))
rect(s5, pt_x5, pt_y5+pt_h5-i(0.25), pt_w5, i(0.25), Y)
txt(s5, "propagandaoca.com.br", pt_x5+i(0.05), pt_y5+pt_h5-i(0.2),
    pt_w5-i(0.1), i(0.18), size=6, color=BK)
for ki in range(8):
    rect(s5, pt_x5+i(0.08), pt_y5+i(0.52)+ki*i(0.33),
         pt_w5-i(0.16), i(0.01), RGBColor(0xD8,0xD8,0xD8))
txt(s5, "PAPEL TIMBRADO A4", pt_x5, pt_y5+pt_h5+i(0.08),
    pt_w5, i(0.2), size=8, color=MG)

# STORY 9:16
st_x5 = pt_x5+pt_w5+i(0.3)
st_w5, st_h5 = i(1.95), i(3.45)
st_y5 = ig_y5
rect(s5, st_x5, st_y5, st_w5, st_h5, Y)
for wi, word in enumerate(["MARKET-","ING QUE","FUNCIO-","NA."]):
    txt(s5, word, st_x5+i(0.14), st_y5+i(0.4)+wi*i(0.64),
        st_w5-i(0.2), i(0.6), size=28, bold=True, color=BK)
rect(s5, st_x5, st_y5+st_h5-i(0.42), st_w5, i(0.42), BK)
txt(s5, "@propaganda.oca", st_x5+i(0.08), st_y5+st_h5-i(0.34),
    st_w5-i(0.1), i(0.26), size=10, color=Y)
txt(s5, "STORY 9:16", st_x5, st_y5+st_h5+i(0.08),
    st_w5, i(0.2), size=8, color=MG)

# EMAIL HEADER mockup
eh_x = st_x5+st_w5+i(0.3)
eh_w = W - eh_x - i(0.3)
eh_y, eh_h = ig_y5, i(1.05)
rect(s5, eh_x, eh_y, eh_w, eh_h, BK)
rect(s5, eh_x, eh_y, i(0.07), eh_h, Y)
add_logo(s5, LOGO_PRINCIPAL, eh_x+i(0.18), eh_y+i(0.08), eh_w-i(0.3), i(0.7))
txt(s5, "oi@propagandaoca.com.br", eh_x+i(0.2), eh_y+i(0.68),
    eh_w-i(0.3), i(0.28), size=10, color=LG)
txt(s5, "ASSINATURA DE E-MAIL", eh_x, ig_y5+eh_h+i(0.08),
    eh_w, i(0.2), size=8, color=MG)

# THUMBNAIL marca
th_y = ig_y5+eh_h+i(0.4)
th_h2 = i(1.9)
rect(s5, eh_x, th_y, int(eh_w/2)-i(0.1), th_h2, Y)
txt(s5, "PO", eh_x+i(0.2), th_y+i(0.3),
    int(eh_w/2)-i(0.3), i(1.2), size=56, bold=True, color=BK)
rect(s5, eh_x+int(eh_w/2)+i(0.1), th_y, int(eh_w/2)-i(0.1), th_h2, BK)
txt(s5, "PO", eh_x+int(eh_w/2)+i(0.3), th_y+i(0.3),
    int(eh_w/2)-i(0.3), i(1.2), size=56, bold=True, color=Y)
txt(s5, "FAVICON  /  AVATAR", eh_x, th_y+th_h2+i(0.08),
    eh_w, i(0.2), size=8, color=MG)

txt(s5, "PROPAGANDA OCA  —  BRANDBOOK 2025",
    i(0.4), i(8.78), i(8), i(0.2), size=8, color=LG)

# ═══════════════════════════════════════════════════════
# SLIDE 6 — TOM DE VOZ
# ═══════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6, OW)

rect(s6, 0, 0, W, i(0.9), Y)
txt(s6, "TOM DE VOZ", i(0.5), i(0.1), i(9), i(0.7),
    size=30, bold=True, color=BK)
txt(s6, "04 / VOZ", i(13.8), i(0.28), i(2), i(0.4),
    size=9, color=BK)

# Coluna esquerda — atributos
attrs = [
    ("DIRETO",    "Vai direto ao ponto. Sem rodeios.\nO empresario nao tem tempo a perder."),
    ("REGIONAL",  "Fala como quem conhece a cidade.\nMenciona lugares e contextos locais."),
    ("CONCRETO",  "Sempre com numeros e resultados\nreais. Nunca promessas vagas."),
    ("PARCEIRO",  "Tom de quem esta do lado.\nSem pedantismo, sem jargao."),
]
ah = int((H - i(0.9) - i(0.2)) / len(attrs))
for ai, (att, desc) in enumerate(attrs):
    ay6 = i(1.0) + ai*ah
    rect(s6, i(0.32), ay6+i(0.08), i(0.055), i(0.38), Y)
    txt(s6, att, i(0.5), ay6, i(6.5), i(0.5),
        size=26, bold=True, color=BK)
    for li, line in enumerate(desc.split("\n")):
        txt(s6, line, i(0.5), ay6+i(0.5)+li*i(0.32),
            i(6.5), i(0.3), size=14, color=MG)

# divisória
rect(s6, i(7.7), i(0.95), i(0.018), H-i(1.1), RGBColor(0xC8,0xC8,0xC8))

# Coluna direita
rx6 = i(7.9)
rw6 = W - rx6 - i(0.4)
txt(s6, "COMO FALAMOS", rx6, i(1.0), rw6, i(0.55),
    size=26, bold=True, color=BK)
rect(s6, rx6, i(1.57), i(1.6), i(0.025), Y)

pares = [
    ("A gente entende o seu mercado.", "Potencialize seus KPIs com\nsolucoes omnichannel."),
    ("Sem enrolacao, sem template.", "Desenvolvemos estrategias\ncustomizadas e escalaveis."),
    ("Resultado em 4 meses.\nNumero real, cliente real.", "ROI otimizado a longo prazo\ncom metodologias ageis."),
    ("Feito para Bom Jesus dos Perdoes.", "Atendemos empresas de todos\nos portes e segmentos."),
]
cw6 = int(rw6/2) - i(0.15)
px_sim = rx6
px_nao = rx6 + cw6 + i(0.18)

py6 = i(1.65)
rect(s6, px_sim, py6, cw6, i(0.32), GBG)
txt(s6, "SIM", px_sim+i(0.12), py6+i(0.04), cw6-i(0.2), i(0.26),
    size=10, bold=True, color=GTX)
rect(s6, px_nao, py6, cw6, i(0.32), RBG)
txt(s6, "NAO", px_nao+i(0.12), py6+i(0.04), cw6-i(0.2), i(0.26),
    size=10, bold=True, color=RTX)
py6 += i(0.34)

bh6 = i(1.15)
for sim, nao in pares:
    rect(s6, px_sim, py6, cw6, bh6, RGBColor(0xF5,0xFC,0xF0))
    for li, line in enumerate(sim.split("\n")):
        txt(s6, line, px_sim+i(0.12), py6+i(0.1)+li*i(0.32),
            cw6-i(0.2), i(0.3), size=13, color=GTX)
    rect(s6, px_nao, py6, cw6, bh6, RGBColor(0xFF,0xF5,0xF5))
    for li, line in enumerate(nao.split("\n")):
        txt(s6, line, px_nao+i(0.12), py6+i(0.1)+li*i(0.32),
            cw6-i(0.2), i(0.3), size=13, color=RTX)
    py6 += bh6 + i(0.05)

# bloco regra de ouro
note_y6 = py6 + i(0.12)
if note_y6 < H - i(0.5):
    note_h6 = H - note_y6 - i(0.22)
    rect(s6, rx6, note_y6, rw6, note_h6, BK)
    txt(s6, "// REGRA DE OURO", rx6+i(0.18), note_y6+i(0.14),
        rw6-i(0.3), i(0.28), size=9, color=Y)
    rect(s6, rx6+i(0.18), note_y6+i(0.46), rw6-i(0.36), i(0.015),
         RGBColor(0x40,0x40,0x40))
    quote = [
        "Se nao daria para dizer para um cliente",
        "olhando nos olhos em Bom Jesus dos Perdoes,",
        "nao coloca no conteudo.",
    ]
    for qi, line in enumerate(quote):
        is_last = (qi == len(quote)-1)
        txt(s6, line, rx6+i(0.18), note_y6+i(0.55)+qi*i(0.34),
            rw6-i(0.3), i(0.32),
            size=13, bold=is_last, color=WH if is_last else LG)

txt(s6, "PROPAGANDA OCA  —  BRANDBOOK 2025",
    i(0.4), i(8.78), i(8), i(0.2), size=8, color=LG)

# ───────────────────────────────────────────────
OUTPUT = "/Users/jessicagregorio/Oca Propaganda/identidade/Identidade_Visual_Editavel.pptx"
prs.save(OUTPUT)
print(f"Salvo: {OUTPUT}  ({len(prs.slides)} slides)")
