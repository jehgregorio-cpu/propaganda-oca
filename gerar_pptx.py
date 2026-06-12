from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

ID_DIR = "/Users/jessicagregorio/Oca Propaganda/identidade"
LOGO_PRINCIPAL = os.path.join(ID_DIR, "logo_pptx_principal.png")
LOGO_AMARELO   = os.path.join(ID_DIR, "logo_pptx_amarelo.png")
LOGO_CLARO     = os.path.join(ID_DIR, "logo_pptx_claro.png")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

YELLOW   = RGBColor(0xFF, 0xDE, 0x21)
BLACK    = RGBColor(0x0A, 0x0A, 0x0A)
OFFWHITE = RGBColor(0xF5, 0xF5, 0xF0)
DGRAY    = RGBColor(0x1A, 0x1A, 0x1A)
MID      = RGBColor(0x44, 0x44, 0x44)
LGRAY    = RGBColor(0xAA, 0xAA, 0xAA)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RED      = RGBColor(0xCC, 0x00, 0x00)
GREEN_BG = RGBColor(0xEA, 0xF3, 0xDE)
GREEN_TX = RGBColor(0x27, 0x50, 0x0A)
RED_BG   = RGBColor(0xFC, 0xEB, 0xEB)
RED_TX   = RGBColor(0xA3, 0x2D, 0x2D)

BLANK = prs.slide_layouts[6]

def add_rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s

def add_text(slide, text, x, y, w, h, size, bold=False, color=BLACK,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Arial"
    return tb

def bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def i(v): return Inches(v)

def add_logo(slide, path, x, y, w, h):
    slide.shapes.add_picture(path, x, y, w, h)

# ==============================================================
# SLIDE 1 — CAPA
# ==============================================================
s1 = prs.slides.add_slide(BLANK)
bg(s1, BLACK)
add_rect(s1, 0, 0, i(0.4), H, YELLOW)
add_logo(s1, LOGO_PRINCIPAL, i(0.7), i(0.9), i(7.5), i(2.8))
add_rect(s1, i(0.7), i(3.75), i(8), i(0.06), YELLOW)
add_text(s1, "Marketing que funciona de verdade.", i(0.7), i(3.7), i(10), i(0.6), 16, color=LGRAY)
add_text(s1, "BRANDBOOK  /  IDENTIDADE VISUAL  /  TOM DE VOZ", i(0.7), i(4.35), i(11), i(0.5), 10, bold=True, color=YELLOW)
add_text(s1, "propagandaoca.com.br  ·  Bom Jesus dos Perdões, SP  ·  2025", i(0.7), i(6.9), i(11), i(0.4), 9, color=RGBColor(0x55,0x55,0x55))

# ==============================================================
# SLIDE 2 — SUMÁRIO
# ==============================================================
s2 = prs.slides.add_slide(BLANK)
bg(s2, OFFWHITE)
add_rect(s2, 0, 0, W, i(1.2), BLACK)
add_text(s2, "SUMÁRIO", i(0.6), i(0.2), i(8), i(0.8), 32, bold=True, color=YELLOW)

items = [
    ("01", "Sobre a Empresa",     "Missão, visão, valores e posicionamento"),
    ("02", "Identidade Visual",   "Cores, tipografia, grid e espaçamento"),
    ("03", "Logotipo",            "Versões, usos corretos e incorretos"),
    ("04", "Tom de Voz",          "Linguagem, exemplos e diretrizes"),
    ("05", "Aplicações","Instagram, cartão de visita, papel timbrado"),
    ("06", "Fotografia e Imagem", "Padrões e diretrizes visuais"),
]
row_h = i(0.82)
for idx, (num, title, sub) in enumerate(items):
    y = i(1.35) + idx * row_h
    add_rect(s2, i(0.5), y, i(0.45), i(0.45), YELLOW)
    add_text(s2, num, i(0.5), y - i(0.03), i(0.45), i(0.5), 10, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s2, title, i(1.05), y, i(4), i(0.3), 13, bold=True, color=BLACK)
    add_text(s2, sub,   i(1.05), y + i(0.28), i(6), i(0.28), 9, color=MID)
    if idx < len(items)-1:
        add_rect(s2, i(0.5), y + i(0.65), i(12.3), i(0.02), RGBColor(0xCC,0xCC,0xCC))

# ==============================================================
# SLIDE 3 — OPENER 01
# ==============================================================
s3 = prs.slides.add_slide(BLANK)
bg(s3, YELLOW)
add_text(s3, "SEÇÃO 01", i(0.7), i(1.8), i(10), i(0.4), 11, bold=True, color=BLACK)
add_text(s3, "SOBRE A EMPRESA", i(0.7), i(2.2), i(12), i(1.4), 52, bold=True, color=BLACK)
add_text(s3, "Quem somos, onde estamos e o que defendemos", i(0.7), i(3.65), i(11), i(0.5), 15, color=BLACK)

# ==============================================================
# SLIDE 4 — MISSÃO VISÃO VALORES
# ==============================================================
s4 = prs.slides.add_slide(BLANK)
bg(s4, OFFWHITE)
add_text(s4, "MISSÃO, VISÃO E VALORES", i(0.6), i(0.3), i(11), i(0.6), 22, bold=True, color=BLACK)
add_rect(s4, i(0.6), i(0.88), i(2.5), i(0.06), YELLOW)

blocks = [
    ("MISSÃO",
     "Ajudar pequenas empresas do interior de SP a crescerem com marketing que faz sentido para o seu contexto. Sem jargão, sem template, sem agência grande que não conhece sua rua."),
    ("VISÃO",
     "Ser a agência de referência em Inbound Marketing e geração de demanda para pequenas empresas de Bom Jesus dos Perdões, Atibaia e Bragança Paulista, reconhecida por resultados concretos e parceria real."),
]
bw = i(5.8); bh = i(2.0)
for idx, (title, body) in enumerate(blocks):
    x = i(0.5) + idx * (bw + i(0.3))
    add_rect(s4, x, i(1.1), bw, bh, BLACK)
    add_text(s4, title, x+i(0.2), i(1.2), bw-i(0.4), i(0.4), 10, bold=True, color=YELLOW)
    add_text(s4, body,  x+i(0.2), i(1.6), bw-i(0.4), i(1.4), 11, color=WHITE)

add_rect(s4, i(0.5), i(3.35), i(12), i(2.8), BLACK)
add_text(s4, "VALORES", i(0.7), i(3.45), i(11), i(0.4), 10, bold=True, color=YELLOW)
valores = [
    ("→ Estratégia local",    "Planos feitos para o seu mercado específico. Nunca templates de agência de São Paulo."),
    ("→ Transparência total", "Relatórios mensais em linguagem humana. Sem jargão, sem enrolação."),
    ("→ Parceria real",            "Sem contratos longos ou multas abusivas. A gente fica porque entrega resultado."),
    ("→ Agilidade",                "Equipe enxuta e focada. Você fala com quem faz, não com assistente de assistente."),
]
vcw = i(2.8)
for idx, (v, desc) in enumerate(valores):
    vx = i(0.6) + idx * (vcw + i(0.18))
    add_text(s4, v,    vx, i(3.85), vcw, i(0.35), 10, bold=True, color=YELLOW)
    add_text(s4, desc, vx, i(4.2),  vcw, i(1.6),  9,  color=LGRAY)

# ==============================================================
# SLIDE 5 — NÚMEROS
# ==============================================================
s5 = prs.slides.add_slide(BLANK)
bg(s5, OFFWHITE)
add_text(s5, "NÚMEROS QUE FALAM POR NÓS", i(0.6), i(0.3), i(11), i(0.6), 22, bold=True, color=BLACK)
add_rect(s5, i(0.6), i(0.88), i(3.5), i(0.06), YELLOW)

stats = [("+120","empresas atendidas"),("3×","mais leads em média"),("87%","taxa de retenção"),("5 anos","no mercado regional")]
cw = i(2.9); ch = i(2.6); gap = i(0.35)
total = 4*cw + 3*gap
start_x = (W - total)/2
for idx, (num, lbl) in enumerate(stats):
    cx = start_x + idx*(cw+gap)
    add_rect(s5, cx, i(1.5), cw, ch, BLACK)
    add_text(s5, num, cx+i(0.2), i(1.7), cw-i(0.3), i(1.3), 40, bold=True, color=YELLOW)
    add_text(s5, lbl.upper(), cx+i(0.2), i(3.2), cw-i(0.3), i(0.5), 9, color=LGRAY)

# ==============================================================
# SLIDE 6 — OPENER 02
# ==============================================================
s6 = prs.slides.add_slide(BLANK)
bg(s6, YELLOW)
add_text(s6, "SEÇÃO 02", i(0.7), i(1.8), i(10), i(0.4), 11, bold=True, color=BLACK)
add_text(s6, "IDENTIDADE VISUAL", i(0.7), i(2.2), i(12), i(1.4), 52, bold=True, color=BLACK)
add_text(s6, "Cores, tipografia, grid e espaçamento", i(0.7), i(3.65), i(11), i(0.5), 15, color=BLACK)

# ==============================================================
# SLIDE 7 — PALETA
# ==============================================================
s7 = prs.slides.add_slide(BLANK)
bg(s7, OFFWHITE)
add_text(s7, "PALETA DE CORES", i(0.6), i(0.3), i(11), i(0.6), 22, bold=True, color=BLACK)
add_rect(s7, i(0.6), i(0.88), i(3), i(0.06), YELLOW)

cores = [
    (YELLOW,   "#FFDE21", "RGB 255, 222, 33",  "CMYK 0 / 13 / 87 / 0",  "Amarelo Primário",
     "Cor de identidade. Destaques, CTAs, títulos sobre fundo escuro.", BLACK),
    (BLACK,    "#0A0A0A", "RGB 10, 10, 10",    "CMYK 0 / 0 / 0 / 96",   "Preto",
     "Fundos principais, texto sobre fundo claro.", WHITE),
    (OFFWHITE, "#F5F5F0", "RGB 245, 245, 240", "CMYK 0 / 0 / 2 / 4",    "Off-white",
     "Fundos secundários, backgrounds de seções claras.", BLACK),
    (DGRAY,    "#1A1A1A", "RGB 26, 26, 26",    "CMYK 0 / 0 / 0 / 90",   "Cinza Escuro",
     "Texto secundário, fundos alternativos.", WHITE),
]
sw = i(3.0); sh = i(1.35); gap7 = i(0.26)
total7 = 4*sw + 3*gap7
sx = (W - total7)/2
for idx, (col, hex_, rgb, cmyk, name, uso, tx) in enumerate(cores):
    cx = sx + idx*(sw+gap7)
    r = add_rect(s7, cx, i(1.1), sw, sh, col)
    if col == OFFWHITE:
        r.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); r.line.width = Pt(0.75)
    add_text(s7, name, cx+i(0.1), i(1.15), sw-i(0.2), i(0.35), 10, bold=True, color=tx)
    add_text(s7, hex_, cx+i(0.1), i(1.48), sw-i(0.2), i(0.25), 8, bold=True, color=tx)
    add_text(s7, rgb,  cx+i(0.1), i(2.55), sw-i(0.2), i(0.25), 8, color=MID)
    add_text(s7, cmyk, cx+i(0.1), i(2.8),  sw-i(0.2), i(0.25), 8, color=MID)
    add_text(s7, uso,  cx+i(0.1), i(3.1),  sw-i(0.2), i(0.6),  8, color=MID)

# ==============================================================
# SLIDE 8 — TIPOGRAFIA
# ==============================================================
s8 = prs.slides.add_slide(BLANK)
bg(s8, OFFWHITE)
add_text(s8, "TIPOGRAFIA", i(0.6), i(0.3), i(11), i(0.6), 22, bold=True, color=BLACK)
add_rect(s8, i(0.6), i(0.88), i(2.5), i(0.06), YELLOW)

fontes = [
    ("Bebas Neue", "TÍTULOS E DISPLAY",
     "Aa Bb Cc 0123", 34,
     "Usada em headings principais, nome da marca e chamadas de alto impacto. Nunca usar em corpo de texto."),
    ("Space Grotesk", "CORPO E SUBTÍTULOS",
     "Aa Bb Cc Dd Ee 0123", 20,
     "Fonte principal de leitura. Usada em parágrafos, subtítulos, legendas e formulários."),
    ("DM Mono", "DADOS E DESTAQUES TÉCNICOS",
     "Aa 01 #FFDE21 ->", 16,
     "Usada para números de destaque, métricas, códigos e labels técnicos."),
]
fw = i(12.2); fh = i(1.55); fy = i(1.15)
for f_name, role, sample, ssize, desc in fontes:
    add_rect(s8, i(0.55), fy, fw, fh, BLACK)
    add_text(s8, role,   i(0.8), fy+i(0.1), i(5), i(0.3), 7, bold=True, color=YELLOW)
    add_text(s8, f_name, i(0.8), fy+i(0.35),i(4), i(0.35), 10, bold=True, color=WHITE)
    add_text(s8, sample, i(4.5), fy+i(0.1), i(7.5), i(0.9), ssize, color=WHITE)
    add_text(s8, desc,   i(0.8), fy+i(0.75),i(3.5), i(0.7), 8, color=LGRAY)
    fy += fh + i(0.18)

# ==============================================================
# SLIDE 9 — OPENER 03
# ==============================================================
s9 = prs.slides.add_slide(BLANK)
bg(s9, YELLOW)
add_text(s9, "SEÇÃO 03", i(0.7), i(1.8), i(10), i(0.4), 11, bold=True, color=BLACK)
add_text(s9, "LOGOTIPO", i(0.7), i(2.2), i(12), i(1.4), 52, bold=True, color=BLACK)
add_text(s9, "Versões, usos corretos e área de proteção", i(0.7), i(3.65), i(11), i(0.5), 15, color=BLACK)

# ==============================================================
# SLIDE 10 — LOGOTIPO VERSÕES
# ==============================================================
s10 = prs.slides.add_slide(BLANK)
bg(s10, OFFWHITE)
add_text(s10, "VERSÕES DO LOGOTIPO", i(0.6), i(0.3), i(11), i(0.5), 22, bold=True, color=BLACK)
add_rect(s10, i(0.6), i(0.82), i(3.5), i(0.06), YELLOW)

versoes = [
    (BLACK,    YELLOW, None,  "Versão Principal",  "Fundos escuros, posts, apresentações."),
    (YELLOW,   BLACK,  None,  "Versão Amarela",    "Impressos, papelaria, eventos."),
    (OFFWHITE, BLACK,  RGBColor(0xCC,0xCC,0xCC), "Versão Clara", "Fundos brancos ou off-white."),
]
vw = i(3.8); vh = i(2.2); vgap = i(0.3)
for idx, (bg_c, fg_c, border, title, desc) in enumerate(versoes):
    vx = i(0.5) + idx*(vw+vgap)
    r = add_rect(s10, vx, i(1.05), vw, vh, bg_c)
    if border:
        r.line.color.rgb = border; r.line.width = Pt(0.75)
    logo_paths = [LOGO_PRINCIPAL, LOGO_AMARELO, LOGO_CLARO]
    add_logo(s10, logo_paths[idx], vx+i(0.1), i(1.15), vw-i(0.2), i(1.7))
    add_text(s10, title, vx, i(3.35), vw, i(0.3), 9, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s10, desc,  vx, i(3.65), vw, i(0.3), 8, color=MID, align=PP_ALIGN.CENTER)

add_text(s10, "USOS INCORRETOS", i(0.6), i(4.1), i(8), i(0.4), 12, bold=True, color=BLACK)
add_rect(s10, i(0.6), i(4.5), i(0.06), i(0.06), RED)
erros = [
    "Não distorcer as proporções do logotipo",
    "Não alterar as cores fora deste manual",
    "Não usar sobre fundos que prejudiquem a legibilidade",
    "Não adicionar sombras, contornos ou efeitos",
    "Não rotacionar o logotipo",
    "Não recriar com outras fontes",
]
half = len(erros) // 2
for idx, e in enumerate(erros):
    col = idx // half; row = idx % half
    ex = i(0.5) + col * i(6.2)
    ey = i(4.3) + row * i(0.45)
    add_rect(s10, ex, ey+i(0.08), i(0.14), i(0.14), RED)
    add_text(s10, e, ex+i(0.25), ey, i(5.8), i(0.4), 9, color=BLACK)

# ==============================================================
# SLIDE 11 — OPENER 04
# ==============================================================
s11 = prs.slides.add_slide(BLANK)
bg(s11, YELLOW)
add_text(s11, "SEÇÃO 04", i(0.7), i(1.8), i(10), i(0.4), 11, bold=True, color=BLACK)
add_text(s11, "TOM DE VOZ", i(0.7), i(2.2), i(12), i(1.4), 52, bold=True, color=BLACK)
add_text(s11, "Como a Propaganda Oca fala com o mundo", i(0.7), i(3.65), i(11), i(0.5), 15, color=BLACK)

# ==============================================================
# SLIDE 12 — TOM DE VOZ
# ==============================================================
s12 = prs.slides.add_slide(BLANK)
bg(s12, OFFWHITE)
add_text(s12, "TOM DE VOZ", i(0.6), i(0.3), i(11), i(0.5), 22, bold=True, color=BLACK)
add_rect(s12, i(0.6), i(0.82), i(2.5), i(0.06), YELLOW)

atts = ["DIRETO", "REGIONAL", "CONCRETO", "PARCEIRO"]
aw = i(2.9)
for idx, att in enumerate(atts):
    ax = i(0.5) + idx*(aw+i(0.28))
    add_rect(s12, ax, i(1.05), aw, i(0.5), YELLOW)
    add_text(s12, att, ax, i(1.08), aw, i(0.44), 13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

pares = [
    ("A gente entende o seu mercado.", "Potencialize seus KPIs com soluções omnichannel."),
    ("Sem enrolação, sem template.", "Desenvolvemos estratégias customizadas e escaláveis."),
    ("Resultado em 4 meses. Número real, cliente real.", "ROI otimizado a longo prazo com metodologias ágeis."),
    ("Feito para Bom Jesus dos Perdões. Feito para você.", "Atendemos empresas de todos os portes e segmentos."),
]
pw = i(5.9); ph12 = i(0.8)
add_text(s12, "SIM",  i(0.5),  i(1.75), pw, i(0.35), 10, bold=True, color=GREEN_TX, align=PP_ALIGN.CENTER)
add_text(s12, "NÃO", i(6.9), i(1.75), pw, i(0.35), 10, bold=True, color=RED_TX, align=PP_ALIGN.CENTER)

for idx, (sim, nao) in enumerate(pares):
    py = i(2.05) + idx*(ph12+i(0.1))
    r1 = add_rect(s12, i(0.5), py, pw, ph12, GREEN_BG)
    add_text(s12, sim, i(0.7), py+i(0.1), pw-i(0.4), ph12-i(0.15), 10, color=GREEN_TX)
    r2 = add_rect(s12, i(6.9), py, pw, ph12, RED_BG)
    add_text(s12, nao, i(7.1), py+i(0.1), pw-i(0.4), ph12-i(0.15), 10, color=RED_TX)

# ==============================================================
# SLIDE 13 — OPENER 05
# ==============================================================
s13 = prs.slides.add_slide(BLANK)
bg(s13, YELLOW)
add_text(s13, "SEÇÃO 05", i(0.7), i(1.8), i(10), i(0.4), 11, bold=True, color=BLACK)
add_text(s13, "APLICAÇÕES", i(0.7), i(2.2), i(12), i(1.4), 52, bold=True, color=BLACK)
add_text(s13, "Instagram, cartão de visita, papel timbrado", i(0.7), i(3.65), i(11), i(0.5), 15, color=BLACK)

# ==============================================================
# SLIDE 14 — APLICAÇÕES
# ==============================================================
s14 = prs.slides.add_slide(BLANK)
bg(s14, OFFWHITE)
add_text(s14, "APLICAÇÕES", i(0.6), i(0.3), i(11), i(0.5), 22, bold=True, color=BLACK)
add_rect(s14, i(0.6), i(0.82), i(2.5), i(0.06), YELLOW)

# Cartão frente
add_text(s14, "CARTÃO DE VISITA", i(0.6), i(1.1), i(5), i(0.35), 11, bold=True, color=BLACK)
cw14 = i(3.8); ch14 = i(2.35)
add_rect(s14, i(0.5), i(1.5), cw14, ch14, BLACK)
add_rect(s14, i(0.5), i(1.5), i(0.22), ch14, YELLOW)
add_logo(s14, LOGO_PRINCIPAL, i(0.8), i(1.6), cw14-i(0.5), i(1.0))
add_text(s14, "Jéssica Gregório", i(0.85), i(2.6),  cw14-i(0.5), i(0.35), 9,  bold=True, color=WHITE)
add_text(s14, "Diretora de Marketing",    i(0.85), i(2.9),  cw14-i(0.5), i(0.25), 7.5, color=LGRAY)
add_text(s14, "(11) 9 5912-5935  ·  oi@propagandaoca.com.br", i(0.85), i(3.6), cw14-i(0.5), i(0.25), 7, color=RGBColor(0x66,0x66,0x66))
add_text(s14, "Frente", i(0.5), i(3.9), cw14, i(0.25), 8, color=MID)

# Cartão verso
vx14 = i(4.5)
add_rect(s14, vx14, i(1.5), cw14, ch14, YELLOW)
add_logo(s14, LOGO_AMARELO, vx14+i(0.2), i(1.6), cw14-i(0.4), i(1.1))
add_text(s14, "Verso", vx14, i(3.9), cw14, i(0.25), 8, color=MID)

# Instagram mockup
add_text(s14, "POST INSTAGRAM", i(8.5), i(1.1), i(4.5), i(0.35), 11, bold=True, color=BLACK)
iw = i(2.9); ih = i(2.9)
add_rect(s14, i(8.5), i(1.5), iw, ih, BLACK)
add_rect(s14, i(8.5), i(1.5)+ih-i(0.45), iw, i(0.45), YELLOW)
add_text(s14, "@propaganda.oca", i(8.55), i(1.52)+ih-i(0.45), iw-i(0.1), i(0.4), 7, bold=True, color=BLACK)
add_text(s14, "3× MAIS LEADS", i(8.6), i(2.3), iw-i(0.2), i(0.8), 20, bold=True, color=YELLOW)
add_text(s14, "em 4 meses de estratégia", i(8.6), i(3.1), iw-i(0.2), i(0.35), 8, color=WHITE)
add_text(s14, "Arquitetura LF · Bom Jesus dos Perdões", i(8.6), i(3.45), iw-i(0.2), i(0.3), 7, color=LGRAY)

# Papel timbrado
add_text(s14, "PAPEL TIMBRADO", i(11.5), i(1.1), i(1.7), i(0.35), 8, bold=True, color=BLACK)
add_rect(s14, i(11.5), i(1.5), i(1.7), i(2.9), WHITE, line=RGBColor(0xCC,0xCC,0xCC))
add_rect(s14, i(11.5), i(1.5)+i(2.9)-i(0.5), i(1.7), i(0.5), BLACK)
add_text(s14, "PROP.", i(11.55), i(1.5)+i(2.9)-i(0.45), i(1.6), i(0.22), 6, bold=True, color=YELLOW)
add_text(s14, "OCA",   i(11.55), i(1.5)+i(2.9)-i(0.25), i(1.6), i(0.22), 6, bold=True, color=YELLOW)
add_rect(s14, i(11.5), i(1.5), i(1.7), i(0.3), YELLOW)
add_text(s14, "propagandaoca.com.br", i(11.5), i(1.5), i(1.7), i(0.28), 5, color=BLACK, align=PP_ALIGN.CENTER)
for k in range(4):
    add_rect(s14, i(11.6), i(2.1)+k*i(0.35), i(1.4), i(0.08), RGBColor(0xDD,0xDD,0xDD))

# ==============================================================
# SLIDE 15 — OPENER 06
# ==============================================================
s15 = prs.slides.add_slide(BLANK)
bg(s15, YELLOW)
add_text(s15, "SEÇÃO 06", i(0.7), i(1.8), i(10), i(0.4), 11, bold=True, color=BLACK)
add_text(s15, "FOTOGRAFIA E IMAGEM", i(0.7), i(2.2), i(12), i(1.4), 46, bold=True, color=BLACK)
add_text(s15, "Padrões visuais e diretrizes de uso de imagem", i(0.7), i(3.65), i(11), i(0.5), 15, color=BLACK)

# ==============================================================
# SLIDE 16 — FOTOGRAFIA
# ==============================================================
s16 = prs.slides.add_slide(BLANK)
bg(s16, OFFWHITE)
add_text(s16, "PADRÕES DE IMAGEM", i(0.6), i(0.3), i(11), i(0.5), 22, bold=True, color=BLACK)
add_rect(s16, i(0.6), i(0.82), i(3.5), i(0.06), YELLOW)

diretrizes = [
    ("Fotografia documental",
     "Imagens reais de estabelecimentos e empresários da região. Sem banco de imagens genérico."),
    ("Tratamento de cor",
     "Contraste alto, tons desaturados. Sobreposição amarela (#FFDE21) a 20-30% de opacidade como recurso de identidade."),
    ("Enquadramento",
     "Composição geométrica. Espaço negativo para texto. Preferência por formato 1:1 ou 4:5 para redes sociais."),
    ("Pessoas",
     "Fotos de clientes e donos de negócio reais em contexto de trabalho. Evitar poses ensaiadas."),
    ("Ícones e ilustrações",
     "Somente ícones outline minimalistas na paleta da marca. Sem clipart ou ilustrações complexas."),
]
dw = i(12.2); col_w = (dw - i(0.3)) / 2
dy = i(1.1)
for idx, (title, desc) in enumerate(diretrizes):
    if idx < 3:
        dx = i(0.55)
        dy_row = i(1.1) + idx * i(1.45)
    else:
        dx = i(0.55) + (idx-3) * (col_w + i(0.3))
        dy_row = i(5.45)
    add_rect(s16, dx, dy_row, col_w if idx >= 3 else dw, i(1.3), BLACK)
    add_text(s16, title, dx+i(0.2), dy_row+i(0.1), (col_w if idx>=3 else dw)-i(0.4), i(0.35), 10, bold=True, color=YELLOW)
    add_text(s16, desc,  dx+i(0.2), dy_row+i(0.45), (col_w if idx>=3 else dw)-i(0.4), i(0.8), 9, color=LGRAY)

# ==============================================================
# SLIDE 17 — CONTRACAPA
# ==============================================================
s17 = prs.slides.add_slide(BLANK)
bg(s17, YELLOW)
add_logo(s17, LOGO_AMARELO, i(2.0), i(0.9), i(9.3), i(2.7))
add_rect(s17, i(3), i(3.9), i(7.3), i(0.06), BLACK)
add_text(s17, "propagandaoca.com.br", i(0.7), i(4.05), W-i(1.4), i(0.5), 14, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s17, "oi@propagandaoca.com.br  ·  (11) 9 5912-5935", i(0.7), i(4.55), W-i(1.4), i(0.4), 12, color=BLACK, align=PP_ALIGN.CENTER)
add_text(s17, "Bom Jesus dos Perdões, SP  ·  2025  ·  Todos os direitos reservados.", i(0.7), i(6.9), W-i(1.4), i(0.4), 9, color=RGBColor(0x33,0x33,0x33), align=PP_ALIGN.CENTER)

OUTPUT = "/Users/jessicagregorio/Oca Propaganda/Brandbook_Propaganda_Oca.pptx"
prs.save(OUTPUT)
print("PPTX salvo:", OUTPUT)
